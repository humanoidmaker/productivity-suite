from __future__ import annotations
"""Export presentation data to PPTX using python-pptx."""

import io

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def export_presentation(slides_meta: dict, theme: dict, title: str) -> bytes:
    """Convert internal presentation data to PPTX bytes."""
    prs = Presentation()

    # Set slide dimensions (16:9 default)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = slides_meta.get("slides", [])
    elements_data = slides_meta.get("elements", {})

    for slide_info in slides_data:
        slide_id = slide_info.get("id", "")
        layout_name = slide_info.get("layout", "blank")

        # Use blank layout for flexibility
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        # Add elements
        slide_elements = elements_data.get(slide_id, [])
        for elem in slide_elements:
            elem_type = elem.get("type", "textbox")
            x = Inches(elem.get("x", 1))
            y = Inches(elem.get("y", 1))
            w = Inches(elem.get("width", 4))
            h = Inches(elem.get("height", 1))

            if elem_type == "textbox":
                txBox = slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True
                text = elem.get("text", "")
                p = tf.paragraphs[0]
                p.text = text
                font_size = elem.get("fontSize", 18)
                p.font.size = Pt(font_size)
                if elem.get("bold"):
                    p.font.bold = True
                if elem.get("color"):
                    try:
                        p.font.color.rgb = RGBColor.from_string(elem["color"].lstrip("#"))
                    except Exception:
                        pass
                align = elem.get("textAlign", "left")
                if align == "center":
                    p.alignment = PP_ALIGN.CENTER
                elif align == "right":
                    p.alignment = PP_ALIGN.RIGHT

            elif elem_type == "shape":
                from pptx.enum.shapes import MSO_SHAPE
                shape_type = elem.get("shape", "RECTANGLE")
                mso = getattr(MSO_SHAPE, shape_type.upper(), MSO_SHAPE.RECTANGLE)
                shape = slide.shapes.add_shape(mso, x, y, w, h)
                if elem.get("fill"):
                    try:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor.from_string(elem["fill"].lstrip("#"))
                    except Exception:
                        pass

            elif elem_type == "image":
                image_path = elem.get("src")
                if image_path:
                    try:
                        slide.shapes.add_picture(image_path, x, y, w, h)
                    except Exception:
                        pass

        # Speaker notes
        notes = slide_info.get("speakerNotes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    prs.core_properties.title = title

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

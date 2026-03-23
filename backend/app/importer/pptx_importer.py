from __future__ import annotations
"""Import PPTX to internal presentation format using python-pptx."""

import io
import uuid

from pptx import Presentation
from pptx.util import Emu


def import_pptx(data: bytes) -> tuple[dict, dict]:
    """Convert PPTX bytes to (slides_meta, theme) dicts."""
    prs = Presentation(io.BytesIO(data))

    slides_info = []
    elements_data = {}

    for slide_idx, slide in enumerate(prs.slides):
        slide_id = f"slide-{uuid.uuid4().hex[:8]}"
        notes_text = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text or ""

        slides_info.append({
            "id": slide_id,
            "layout": "blank",
            "transition": "none",
            "transitionDuration": 0.5,
            "speakerNotes": notes_text,
            "hidden": False,
        })

        slide_elements = []
        for shape in slide.shapes:
            elem: dict = {
                "type": "textbox",
                "x": _emu_to_inches(shape.left),
                "y": _emu_to_inches(shape.top),
                "width": _emu_to_inches(shape.width),
                "height": _emu_to_inches(shape.height),
            }

            if shape.has_text_frame:
                text_parts = []
                for para in shape.text_frame.paragraphs:
                    text_parts.append(para.text)
                elem["text"] = "\n".join(text_parts)
                elem["type"] = "textbox"
                # Get font info from first run
                if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                    run = shape.text_frame.paragraphs[0].runs[0]
                    if run.font.size:
                        elem["fontSize"] = round(run.font.size.pt)
                    if run.font.bold:
                        elem["bold"] = True
            elif shape.shape_type and "PICTURE" in str(shape.shape_type):
                elem["type"] = "image"
                elem["src"] = None  # Image extraction requires saving to file
            else:
                elem["type"] = "shape"
                elem["shape"] = "RECTANGLE"

            slide_elements.append(elem)

        elements_data[slide_id] = slide_elements

    slides_meta = {"slides": slides_info, "elements": elements_data}

    # Extract theme colors (simplified)
    theme = {
        "primary": "#2563eb",
        "secondary": "#64748b",
        "accent": "#f59e0b",
        "background": "#ffffff",
        "text": "#1e293b",
        "headingFont": "Calibri",
        "bodyFont": "Calibri",
    }

    return slides_meta, theme


def _emu_to_inches(emu: int | None) -> float:
    if emu is None:
        return 0
    return round(emu / 914400, 2)
